"""
AgriTrust — Project Review Presentation Generator
Generates a professional PPT with light background and formal tone.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors ───
WHITE       = RGBColor(255, 255, 255)
LIGHT_BG    = RGBColor(245, 248, 245)  # Very light green-grey
DARK_TEXT    = RGBColor(30, 40, 35)
GREEN_PRIMARY = RGBColor(27, 94, 32)    # #1B5E20
GREEN_ACCENT  = RGBColor(67, 160, 71)   # #43A047
GREEN_LIGHT   = RGBColor(200, 230, 201) # #C8E6C9
GREY_TEXT     = RGBColor(80, 90, 85)
GREY_LIGHT    = RGBColor(220, 225, 220)
ORANGE        = RGBColor(230, 126, 34)
RED_SOFT      = RGBColor(192, 57, 43)
BLUE_SOFT     = RGBColor(41, 128, 185)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════

def add_bg(slide, color=LIGHT_BG):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with single-style text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=DARK_TEXT, spacing=Pt(6), bullet_char="●"):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox

def add_top_bar(slide):
    """Add a thin green accent bar at the top."""
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), GREEN_PRIMARY)

def add_bottom_bar(slide, slide_num, total=16):
    """Add footer with slide number."""
    add_rect(slide, Inches(0), Inches(7.18), prs.slide_width, Inches(0.32), GREEN_PRIMARY)
    add_text_box(slide, Inches(0.5), Inches(7.18), Inches(5), Inches(0.32),
                 "AgriTrust — Project Review Presentation", font_size=9,
                 color=WHITE, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(10), Inches(7.18), Inches(3), Inches(0.32),
                 f"{slide_num} / {total}", font_size=9,
                 color=WHITE, alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, section_num, title, subtitle=""):
    """Add a section divider slide."""
    add_bg(slide, GREEN_PRIMARY)
    # Section number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(1.8), Inches(1.7), Inches(1.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = WHITE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = GREEN_PRIMARY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(18)

    add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
                 title, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                     subtitle, font_size=16, color=RGBColor(200, 230, 201), alignment=PP_ALIGN.CENTER)

def add_slide_title(slide, title, subtitle=""):
    """Add title + green underline to a content slide."""
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
                 title, font_size=26, bold=True, color=GREEN_PRIMARY)
    add_rect(slide, Inches(0.8), Inches(0.92), Inches(2.5), Inches(0.05), GREEN_ACCENT)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
                     subtitle, font_size=13, color=GREY_TEXT)

def make_table_row_card(slide, left, top, width, col1, col2, fill, text_color=DARK_TEXT):
    """Simple 2-column row card."""
    card = add_rounded_rect(slide, left, top, width, Inches(0.45), fill, border_color=GREY_LIGHT)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.05), Inches(2.2), Inches(0.35),
                 col1, font_size=12, bold=True, color=text_color)
    add_text_box(slide, left + Inches(2.5), top + Inches(0.05), width - Inches(2.7), Inches(0.35),
                 col2, font_size=12, color=text_color)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(3.8), GREEN_PRIMARY)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.5),
             "🌾", font_size=40, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
             "AgriTrust", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10), Inches(0.6),
             "Empowering Small Farmers with AI-Driven Government Scheme Discovery",
             font_size=18, color=RGBColor(200, 230, 201), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.5),
             "Second Review — Progress & Roadmap", font_size=14,
             color=RGBColor(180, 210, 185), alignment=PP_ALIGN.CENTER)

# Info below green section
add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.4),
             "Technology Stack:  Flutter (Mobile)  •  Python Flask (Backend)  •  MongoDB (Database)  •  Gemini AI",
             font_size=13, color=GREY_TEXT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.4),
             "Team Jeeva  |  February 2026",
             font_size=14, bold=True, color=GREEN_PRIMARY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.18), prs.slide_width, Inches(0.32), GREEN_PRIMARY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Agenda")

agenda_items = [
    ("1", "Problem Refinement", "Panel feedback incorporation & target user finalization"),
    ("2", "Implementation Progress", "Core features developed and running on localhost"),
    ("3", "Challenges Faced", "Technical issues, scope adjustments & time constraints"),
    ("4", "Scaling Plan", "Cloud hosting, deployment & real-time data strategies"),
    ("5", "Timeline till Hackathon End", "Milestone-based development roadmap"),
    ("6", "Novelty", "AI/ML integration & unique differentiators"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.5) + Inches(i * 0.85)
    card = add_rounded_rect(slide, Inches(1.5), y, Inches(10), Inches(0.7), WHITE, border_color=GREEN_LIGHT)
    
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.7), y + Inches(0.1), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = GREEN_PRIMARY
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.4), y + Inches(0.05), Inches(4), Inches(0.35),
                 title, font_size=16, bold=True, color=GREEN_PRIMARY)
    add_text_box(slide, Inches(2.4), y + Inches(0.38), Inches(8.5), Inches(0.3),
                 desc, font_size=11, color=GREY_TEXT)

add_bottom_bar(slide, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — SECTION 1 DIVIDER: PROBLEM REFINEMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 1, "Problem Refinement",
                   "Panel Feedback Incorporation & Target User Finalization")

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — PANEL FEEDBACK INCORPORATED
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Panel Feedback — Key Refinements")

feedbacks = [
    ("Scope Narrowed",
     "Shifted from \"all farmers\" to small & marginal farmers (0.5–10 hectares) who are digital-first-time users, unfamiliar with government portals."),
    ("Language Barrier Addressed",
     "English-only was an exclusion factor. App now supports multilingual access — English & Tamil — with architecture ready for Hindi and Malayalam."),
    ("Voice-First Interaction",
     "Integrated Text-to-Speech (TTS) engine so scheme explanations are read aloud in the farmer's chosen language, addressing low literacy."),
    ("Feature Scope Reduced",
     "Dropped complex admin and advisory modules. Finalized to three core needs: scheme discovery, weather forecasts, and crop market prices."),
]

for i, (title, desc) in enumerate(feedbacks):
    y = Inches(1.4) + Inches(i * 1.4)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.2),
                             RGBColor(248, 252, 248), border_color=GREEN_LIGHT)
    
    # Green left accent
    add_rect(slide, Inches(0.8), y, Inches(0.07), Inches(1.2), GREEN_ACCENT)
    
    add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(10.5), Inches(0.3),
                 f"✅  {title}", font_size=15, bold=True, color=GREEN_PRIMARY)
    add_text_box(slide, Inches(1.3), y + Inches(0.48), Inches(10.5), Inches(0.65),
                 desc, font_size=12, color=GREY_TEXT)

add_bottom_bar(slide, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — TARGET USER GROUP
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Finalized Target User Group")

rows = [
    ("Who", "Small & marginal farmers (0.5–10 hectares landholding)"),
    ("Where", "Rural India — prioritizing Tamil Nadu initially"),
    ("Language", "Low-literacy users; prefer local language + voice interaction"),
    ("Device", "Android mobile (primary), with web/desktop fallback"),
    ("Context", "Limited internet; seeking government scheme benefits"),
    ("Pain Point", "Unaware of eligible schemes; can't navigate complex portals"),
]

# Table header
add_rounded_rect(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.55),
                 GREEN_PRIMARY)
add_text_box(slide, Inches(1.7), Inches(1.55), Inches(2.2), Inches(0.45),
             "Attribute", font_size=14, bold=True, color=WHITE)
add_text_box(slide, Inches(4.0), Inches(1.55), Inches(7.2), Inches(0.45),
             "Definition", font_size=14, bold=True, color=WHITE)

for i, (attr, defn) in enumerate(rows):
    y = Inches(2.15) + Inches(i * 0.62)
    fill = RGBColor(248, 252, 248) if i % 2 == 0 else WHITE
    add_rounded_rect(slide, Inches(1.5), y, Inches(10), Inches(0.55), fill, border_color=GREY_LIGHT)
    add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(2.2), Inches(0.4),
                 attr, font_size=13, bold=True, color=GREEN_PRIMARY)
    add_text_box(slide, Inches(4.0), y + Inches(0.08), Inches(7.2), Inches(0.4),
                 defn, font_size=12, color=DARK_TEXT)

add_bottom_bar(slide, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — SECTION 2 DIVIDER: IMPLEMENTATION PROGRESS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 2, "Implementation Progress",
                   "Core Features Developed — Running Successfully on Localhost")

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "System Architecture Overview")

# Architecture boxes
boxes = [
    (Inches(0.6), Inches(2.0), Inches(3.6), Inches(4.2),
     "📱 Flutter Mobile App", [
         "Landing Screen — Language selection",
         "Farmer Input — Crop, land, season, state",
         "Scheme Recommendations — Sorted list",
         "Scheme Detail — TTS voice explanation",
         "Dashboard — Weather + Market Prices",
         "",
         "Services:",
         "  • TTS / STT (voice I/O)",
         "  • Language (EN / Tamil)",
         "  • API Service (HTTP client)",
         "  • Local Storage (Hive cache)",
     ]),
    (Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.2),
     "⚙️ Flask Backend (Python)", [
         "API Routes:",
         "  POST /api/getEligibleSchemes",
         "  GET  /api/weather",
         "  GET  /api/market-prices",
         "  POST /api/addScheme",
         "",
         "Services:",
         "  • Eligibility Matching Engine",
         "  • Weather Service (Open-Meteo)",
         "  • Market Price Service (MSP)",
         "  • Web Scraper (5 gov sources)",
         "  • PDF Parser",
     ]),
    (Inches(9.0), Inches(2.0), Inches(3.6), Inches(4.2),
     "🗄️ Data Layer", [
         "MongoDB:",
         "  • Schemes collection",
         "  • Eligibility criteria",
         "  • Localized descriptions",
         "",
         "External APIs:",
         "  • Open-Meteo (weather)",
         "  • Geocoding (GPS → state)",
         "",
         "Seed Data:",
         "  • 8 major national schemes",
         "  • 23+ crop MSP prices",
     ]),
]

for (left, top, w, h, title, items) in boxes:
    card = add_rounded_rect(slide, left, top, w, h,
                             RGBColor(250, 253, 250), border_color=GREEN_LIGHT)
    add_rect(slide, left, top, w, Inches(0.06), GREEN_ACCENT)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), w - Inches(0.3), Inches(0.35),
                 title, font_size=14, bold=True, color=GREEN_PRIMARY)
    
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.55),
                                     w - Inches(0.3), h - Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(10)
        p.font.color.rgb = GREY_TEXT
        p.font.name = "Calibri"
        p.space_after = Pt(2)

# Arrows between boxes
for x in [Inches(4.2), Inches(8.6)]:
    add_text_box(slide, x, Inches(3.8), Inches(0.6), Inches(0.5),
                 "→", font_size=28, bold=True, color=GREEN_ACCENT, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — FRONTEND SCREENS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Flutter Frontend — Screens Developed")

screens = [
    ("Landing Screen", "✅ Done", "Language selection (EN / Tamil), onboarding flow with smooth animations"),
    ("Farmer Input Screen", "✅ Done", "Crop type, land size stepper, season chips, searchable state picker with GPS auto-detect"),
    ("Scheme Recommendations", "✅ Done", "Displays matched government schemes sorted by benefit amount with voice read-aloud"),
    ("Scheme Detail Screen", "✅ Done", "Full scheme info — documents checklist, eligibility criteria, official link, TTS explanation"),
    ("Dashboard Screen", "✅ Done", "Weather tab (current + 5-day forecast) and live crop market prices by state with trends"),
]

# Header
add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5), GREEN_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(1.44), Inches(3), Inches(0.42),
             "Screen", font_size=13, bold=True, color=WHITE)
add_text_box(slide, Inches(3.8), Inches(1.44), Inches(1.2), Inches(0.42),
             "Status", font_size=13, bold=True, color=WHITE)
add_text_box(slide, Inches(5.2), Inches(1.44), Inches(7), Inches(0.42),
             "Description", font_size=13, bold=True, color=WHITE)

for i, (screen, status, desc) in enumerate(screens):
    y = Inches(2.0) + Inches(i * 1.0)
    fill = RGBColor(248, 252, 248) if i % 2 == 0 else WHITE
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.85), fill, border_color=GREY_LIGHT)
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(2.8), Inches(0.55),
                 screen, font_size=13, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(3.8), y + Inches(0.15), Inches(1.2), Inches(0.55),
                 status, font_size=12, bold=True, color=GREEN_ACCENT)
    add_text_box(slide, Inches(5.2), y + Inches(0.15), Inches(7), Inches(0.55),
                 desc, font_size=11, color=GREY_TEXT)

add_bottom_bar(slide, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — BACKEND MODULES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Backend Modules — Flask + MongoDB")

modules = [
    ("Eligibility Engine", "/api/getEligibleSchemes", "Rule-based filter: state, crop, land size, season → MongoDB query → sorted by benefit"),
    ("Weather Service", "/api/weather", "Open-Meteo free API → current conditions + 5-day forecast with WMO weather code decoding"),
    ("Market Prices", "/api/market-prices", "MSP-based prices with daily fluctuation for 23+ crops across 14 states and key mandis"),
    ("Database Seeding", "seed_db.py", "MongoDB seeded with PM-KISAN, PMFBY, KCC, Soil Health Card, NMSA, PKVY, RKVY, Micro Irrigation"),
    ("Web Scraper", "scraper.py", "Scrapes 5 government sources: data.gov.in, agricoop, pmkisan, pmfby, nabard"),
    ("API Security", "Input validation", "Sanitization, NoSQL injection prevention, pagination, coordinate validation"),
]

# Header
add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5), GREEN_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(1.44), Inches(2.5), Inches(0.42),
             "Module", font_size=13, bold=True, color=WHITE)
add_text_box(slide, Inches(3.5), Inches(1.44), Inches(2.2), Inches(0.42),
             "Endpoint / File", font_size=13, bold=True, color=WHITE)
add_text_box(slide, Inches(5.8), Inches(1.44), Inches(6.5), Inches(0.42),
             "Description", font_size=13, bold=True, color=WHITE)

for i, (module, endpoint, desc) in enumerate(modules):
    y = Inches(2.0) + Inches(i * 0.82)
    fill = RGBColor(248, 252, 248) if i % 2 == 0 else WHITE
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(0.72), fill, border_color=GREY_LIGHT)
    add_text_box(slide, Inches(0.8), y + Inches(0.12), Inches(2.5), Inches(0.5),
                 module, font_size=12, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(3.5), y + Inches(0.12), Inches(2.2), Inches(0.5),
                 endpoint, font_size=10, color=BLUE_SOFT)
    add_text_box(slide, Inches(5.8), y + Inches(0.12), Inches(6.5), Inches(0.5),
                 desc, font_size=11, color=GREY_TEXT)

add_bottom_bar(slide, 9)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — SECTION 3 DIVIDER: CHALLENGES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 3, "Challenges Faced",
                   "Technical Integration Issues, Scope Adjustments & Time Constraints")

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — TECHNICAL CHALLENGES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Technical Integration Challenges & Resolutions")

challenges = [
    ("⚠️  Mobile-to-Backend Connectivity",
     "App called localhost — unreachable from physical Android device.",
     "Configured API base URL to machine's LAN IP (192.168.x.x:5000) with dynamic IP config — no rebuild needed."),
    ("⚠️  Location & GPS Permissions",
     "Android requires explicit runtime permissions; user may permanently deny.",
     "Implemented graceful fallback — if GPS denied, defaults to Chennai, Tamil Nadu."),
    ("⚠️  Speech-to-Text on Web vs Mobile",
     "speech_to_text Flutter package has different behavior on web vs Android.",
     "Platform-detection logic: STT enabled on mobile; web users prompted for text entry."),
    ("⚠️  Empty MongoDB Database",
     "Seed file was empty — no schemes visible to users on first load.",
     "Populated seed_db.py with real government scheme data including eligibility, benefits, and documents."),
]

for i, (title, problem, resolution) in enumerate(challenges):
    y = Inches(1.4) + Inches(i * 1.4)
    card = add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1.25),
                             RGBColor(255, 252, 248), border_color=RGBColor(255, 224, 178))
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(11.5), Inches(0.3),
                 title, font_size=14, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(1.0), y + Inches(0.38), Inches(11), Inches(0.3),
                 f"Problem: {problem}", font_size=11, color=RED_SOFT)
    add_text_box(slide, Inches(1.0), y + Inches(0.68), Inches(11), Inches(0.5),
                 f"✅ Resolution: {resolution}", font_size=11, color=GREEN_PRIMARY)

add_bottom_bar(slide, 11)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — SCOPE ADJUSTMENTS + TIME
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Scope Adjustments & Time Constraints")

# Left column — Scope
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2),
                 RGBColor(250, 253, 250), border_color=GREEN_LIGHT)
add_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.06), GREEN_ACCENT)
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5), Inches(0.35),
             "📏  Scope Adjustments", font_size=16, bold=True, color=GREEN_PRIMARY)

scope_items = [
    "Removed live market scraping from Agmarknet — unreliable due to CAPTCHA and frequent HTML structure changes.",
    "Replaced with deterministic MSP-based simulation using official 2024–25 prices with time-seeded daily fluctuation.",
    "Removed admin dashboard for scheme management — deprioritized to focus on core farmer-facing experience.",
]
add_bullet_list(slide, Inches(0.9), Inches(2.3), Inches(5.2), Inches(3.5),
                scope_items, font_size=12, color=GREY_TEXT, bullet_char="→")

# Right column — Time
add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2),
                 RGBColor(250, 253, 250), border_color=GREEN_LIGHT)
add_rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.06), ORANGE)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.35),
             "⏳  Time Constraints", font_size=16, bold=True, color=ORANGE)

time_items = [
    "Hindi and Malayalam UI strings prepared in localization architecture but not fully wired — Tamil is fully functional.",
    "Unit tests (test_api.py, test_filtering.py) written and passing for core engine. End-to-end mobile testing partially done.",
    "PDF scheme parser developed but importing is manual — automated ingestion pipeline planned but not completed.",
]
add_bullet_list(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(3.5),
                time_items, font_size=12, color=GREY_TEXT, bullet_char="→")

add_bottom_bar(slide, 12)

# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — SECTION 4 DIVIDER: SCALING PLAN
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 4, "Scaling Plan",
                   "Cloud Hosting, Deployment & Real-Time Data Strategies")

# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — SCALING PLAN CONTENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Cloud Hosting & Deployment Strategy")

# 3 columns
cols = [
    (Inches(0.5), "☁️ Backend Deployment", [
        "Deploy Flask API to Railway / Render (free tier)",
        "MongoDB Atlas (cloud) — free 512 MB cluster",
        "Environment variables via platform secrets",
        "Auto-restart + health check monitoring",
        "HTTPS via platform SSL certificates",
    ]),
    (Inches(4.6), "📱 App Distribution", [
        "Generate signed APK for Android",
        "Distribute via Google Play (internal testing)",
        "Web version on Firebase Hosting",
        "ADB sideload for demo devices",
        "OTA updates via CodePush (planned)",
    ]),
    (Inches(8.7), "🔄 Real-Time Data", [
        "Scheduler: Auto-refresh schemes every 24h",
        "data.gov.in API for live Agmarknet prices",
        "Web scraper upgraded with fallback data",
        "Weather: Already real-time (Open-Meteo)",
        "Caching layer for offline-first support",
    ]),
]

for (left, title, items) in cols:
    card = add_rounded_rect(slide, left, Inches(1.5), Inches(3.8), Inches(4.5),
                             RGBColor(250, 253, 250), border_color=GREEN_LIGHT)
    add_rect(slide, left, Inches(1.5), Inches(3.8), Inches(0.06), GREEN_ACCENT)
    add_text_box(slide, left + Inches(0.15), Inches(1.7), Inches(3.5), Inches(0.35),
                 title, font_size=14, bold=True, color=GREEN_PRIMARY)
    add_bullet_list(slide, left + Inches(0.15), Inches(2.2), Inches(3.5), Inches(3.5),
                    items, font_size=11, color=GREY_TEXT, bullet_char="•")

add_bottom_bar(slide, 14)

# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — SECTION 5 DIVIDER: TIMELINE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 5, "Timeline till Hackathon End",
                   "Milestone-Based Development Roadmap")

# ═══════════════════════════════════════════════════════════════
# SLIDE 16 — TIMELINE CONTENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "Development Timeline — Remaining Milestones")

milestones = [
    ("Week 1", "Feb 21 – Feb 28", "AI Integration (Layer 1)",
     "Gemini chatbot endpoint • NLP voice parsing • Smart scheme ranking",
     GREEN_ACCENT),
    ("Week 2", "Mar 1 – Mar 7", "ML Models & Real-Time Data",
     "Crop yield prediction (RandomForest) • Market price API (data.gov.in) • Scheduler setup",
     BLUE_SOFT),
    ("Week 3", "Mar 8 – Mar 14", "Advanced Features",
     "Crop disease detection (MobileNetV3) • Soil health assessment • Hindi localization",
     ORANGE),
    ("Week 4", "Mar 15 – Mar 21", "Deployment & Testing",
     "Cloud deploy (Railway + Atlas) • APK build • End-to-end testing • Demo preparation",
     RGBColor(156, 39, 176)),
]

# Timeline line
add_rect(slide, Inches(1.5), Inches(1.9), Inches(0.04), Inches(5.0), GREEN_LIGHT)

for i, (week, dates, phase, details, color) in enumerate(milestones):
    y = Inches(1.8) + Inches(i * 1.25)
    
    # Timeline dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.35), y + Inches(0.15), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    tf = dot.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Content card
    card = add_rounded_rect(slide, Inches(2.2), y, Inches(10.2), Inches(1.1),
                             WHITE, border_color=GREY_LIGHT)
    add_rect(slide, Inches(2.2), y, Inches(0.06), Inches(1.1), color)
    
    add_text_box(slide, Inches(2.5), y + Inches(0.05), Inches(2), Inches(0.3),
                 week, font_size=14, bold=True, color=color)
    add_text_box(slide, Inches(4.5), y + Inches(0.08), Inches(2.5), Inches(0.25),
                 dates, font_size=10, color=GREY_TEXT)
    add_text_box(slide, Inches(2.5), y + Inches(0.35), Inches(9.5), Inches(0.3),
                 phase, font_size=13, bold=True, color=DARK_TEXT)
    add_text_box(slide, Inches(2.5), y + Inches(0.65), Inches(9.5), Inches(0.4),
                 details, font_size=11, color=GREY_TEXT)

add_bottom_bar(slide, 16)

# ═══════════════════════════════════════════════════════════════
# SLIDE 17 — SECTION 6 DIVIDER: NOVELTY
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, 6, "Novelty",
                   "AI/ML Integration & Unique Differentiators")

# ═══════════════════════════════════════════════════════════════
# SLIDE 18 — NOVELTY: AI/ML FEATURES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "AI & ML — Planned Integration Layers")

layers = [
    ("Layer 1 — Quick AI Wins", GREEN_ACCENT, [
        ("🤖 Gemini Chatbot", "Scheme-context aware conversational AI — explains schemes in farmer's language"),
        ("🗣️ NLP Voice Parsing", "Farmer speaks naturally; Gemini extracts crop, state, land size, season from speech"),
    ]),
    ("Layer 2 — ML Models", BLUE_SOFT, [
        ("📊 Crop Yield Prediction", "RandomForest model trained on ICRISAT/data.gov.in datasets — predicts tonnes/hectare"),
        ("📈 Market Price Forecast", "Facebook Prophet on historical mandi data — 7-day crop price trend prediction"),
        ("🏆 Smart Scheme Ranking", "TF-IDF + Cosine Similarity — personalized relevance scoring for each farmer profile"),
    ]),
    ("Layer 3 — Advanced Features", ORANGE, [
        ("🌿 Crop Disease Detection", "MobileNetV3 (transfer learning) on PlantVillage dataset — 38 disease classes, 93% accuracy"),
        ("🌱 Soil Health Assessment", "XGBoost on Soil Health Card database — district-level NPK prediction + fertilizer advice"),
    ]),
]

y_offset = Inches(1.4)
for layer_title, color, features in layers:
    # Layer header
    add_rounded_rect(slide, Inches(0.6), y_offset, Inches(12), Inches(0.42), color)
    add_text_box(slide, Inches(0.8), y_offset + Inches(0.04), Inches(11), Inches(0.35),
                 layer_title, font_size=13, bold=True, color=WHITE)
    y_offset += Inches(0.5)
    
    for feat_name, feat_desc in features:
        add_rounded_rect(slide, Inches(0.8), y_offset, Inches(11.6), Inches(0.5),
                          RGBColor(250, 253, 250), border_color=GREY_LIGHT)
        add_text_box(slide, Inches(1.0), y_offset + Inches(0.06), Inches(3.2), Inches(0.38),
                     feat_name, font_size=12, bold=True, color=DARK_TEXT)
        add_text_box(slide, Inches(4.2), y_offset + Inches(0.06), Inches(8), Inches(0.38),
                     feat_desc, font_size=11, color=GREY_TEXT)
        y_offset += Inches(0.55)
    
    y_offset += Inches(0.15)

add_bottom_bar(slide, 18)

# ═══════════════════════════════════════════════════════════════
# SLIDE 19 — NOVELTY: WHAT MAKES US DIFFERENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_top_bar(slide)
add_slide_title(slide, "What Makes AgriTrust Novel?")

differentiators = [
    ("🎯  Hyper-Focused User Group",
     "Unlike generic portals, AgriTrust is built specifically for low-literacy small farmers — with voice-first, local-language interaction."),
    ("🧠  AI-Powered, Not Just a Database",
     "Eligibility matching evolves from static rules to ML-ranked personalized recommendations. Gemini chatbot provides conversational scheme help."),
    ("🌐  Offline-First Architecture",
     "Scheme results cached locally via Hive. Farmers in low-connectivity areas can still access previously loaded data without internet."),
    ("🗣️  Voice I/O in Local Languages",
     "Combines STT (speech input) + Gemini NLP (extraction) + TTS (voice output) for a fully voice-driven experience in Tamil and English."),
    ("📊  Actionable Intelligence",
     "Not just scheme info — crop yield prediction, disease detection, market price forecasting, and soil health assessment give farmers real decision-making power."),
]

for i, (title, desc) in enumerate(differentiators):
    y = Inches(1.5) + Inches(i * 1.1)
    card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.95),
                             RGBColor(248, 252, 248), border_color=GREEN_LIGHT)
    add_rect(slide, Inches(0.8), y, Inches(0.07), Inches(0.95), GREEN_ACCENT)
    add_text_box(slide, Inches(1.1), y + Inches(0.08), Inches(10.8), Inches(0.3),
                 title, font_size=14, bold=True, color=GREEN_PRIMARY)
    add_text_box(slide, Inches(1.3), y + Inches(0.42), Inches(10.5), Inches(0.5),
                 desc, font_size=12, color=GREY_TEXT)

add_bottom_bar(slide, 19)

# ═══════════════════════════════════════════════════════════════
# SLIDE 20 — THANK YOU
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_PRIMARY)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "🌾", font_size=48, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1),
             "Thank You", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.5),
             "AgriTrust — Empowering Small Farmers with Technology",
             font_size=18, color=RGBColor(200, 230, 201), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9), Inches(0.4),
             "Questions & Feedback Welcome", font_size=16,
             color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(0.4),
             "Team Jeeva  |  February 2026",
             font_size=14, bold=False, color=RGBColor(180, 210, 185), alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.18), prs.slide_width, Inches(0.32), RGBColor(21, 71, 25))

# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AgriTrust_Review_Presentation.pptx")
prs.save(output_path)
print(f"\n[OK] Presentation saved to: {output_path}")
print(f"     Total slides: {len(prs.slides)}")
